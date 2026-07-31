# -*- coding: utf-8 -*-
"""Gera a apresentação de defesa em PowerPoint (.pptx), editável.

Saída: tex/apresentacao/apresentacao_tcc.pptx  (16:9, 22 slides)

Estrutura: problema → pergunta → dados/método → resultados → validação real
→ CONSIDERAÇÕES FINAIS (objetivos, contribuições, limitações, conclusão)
→ demonstração do app. O slide 22 é reserva (não apresentar).

As figuras vêm de tex/apresentacao/figuras (gerar_figuras_slides.py e
gerar_figura_previsto_real.py). Paleta validada para daltonismo (Okabe-Ito).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

RAIZ = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..')
DIR_FIG = os.path.join(RAIZ, 'tex', 'apresentacao', 'figuras')
SAIDA = os.path.join(RAIZ, 'tex', 'apresentacao', 'apresentacao_tcc.pptx')

# ── Paleta ──────────────────────────────────────────────────────────────────
AZUL = RGBColor(0x00, 0x72, 0xB2)
VERM = RGBColor(0xD5, 0x5E, 0x00)
VERDE = RGBColor(0x00, 0x9E, 0x73)
ROXO = RGBColor(0xCC, 0x79, 0xA7)
TINTA = RGBColor(0x26, 0x2B, 0x30)
TINTA2 = RGBColor(0x5A, 0x64, 0x6E)
FUNDOC = RGBColor(0xED, 0xF1, 0xF4)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_L = RGBColor(0xC8, 0xCE, 0xD4)

FONTE = 'Calibri'
TOTAL = 23          # slides no arquivo (o ultimo e reserva)

# ── Geometria (polegadas) ───────────────────────────────────────────────────
LARG, ALT = 13.333, 7.5
MG = 0.62                    # margem lateral
TOP_TIT = 0.34
TOP_SUB = 0.94
Y_LINHA = 1.42
TOP_CONT = 1.62              # início da área de conteúdo
BASE_CONT = 6.86             # fim da área de conteúdo
Y_RODAPE = 6.95

prs = Presentation()
prs.slide_width = Inches(LARG)
prs.slide_height = Inches(ALT)
BRANCO_LAYOUT = prs.slide_layouts[6]      # em branco


# ── Utilidades ──────────────────────────────────────────────────────────────
def caixa(slide, x, y, w, h, texto, tam=16, cor=TINTA, negrito=False,
          italico=False, alinha=PP_ALIGN.LEFT, ancora=MSO_ANCHOR.TOP,
          entrelinha=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = ancora
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = alinha
    p.line_spacing = entrelinha
    r = p.add_run()
    r.text = texto
    f = r.font
    f.name, f.size, f.color.rgb, f.bold, f.italic = FONTE, Pt(tam), cor, negrito, italico
    return tb


def rico(slide, x, y, w, h, linhas, tam=16, entrelinha=1.2, espaco_antes=6,
         alinha=PP_ALIGN.LEFT, ancora=MSO_ANCHOR.TOP):
    """linhas = lista de parágrafos; cada parágrafo = lista de (texto, dict).

    dict aceita: cor, negrito, italico, tam, fonte.
    """
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = ancora
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, par in enumerate(linhas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alinha
        p.line_spacing = entrelinha
        if i > 0:
            p.space_before = Pt(espaco_antes)
        for txt, est in par:
            r = p.add_run()
            r.text = txt
            f = r.font
            f.name = est.get('fonte', FONTE)
            f.size = Pt(est.get('tam', tam))
            f.color.rgb = est.get('cor', TINTA)
            f.bold = est.get('negrito', False)
            f.italic = est.get('italico', False)
    return tb


def marcadores(slide, x, y, w, h, itens, tam=16, cor_marca=AZUL, entrelinha=1.2,
               espaco=9):
    """itens = lista de parágrafos, cada um lista de (texto, dict)."""
    pars = []
    for it in itens:
        pars.append([('▪  ', {'cor': cor_marca, 'negrito': True, 'tam': tam})] + it)
    return rico(slide, x, y, w, h, pars, tam=tam, entrelinha=entrelinha,
                espaco_antes=espaco)


def titulo(slide, texto, sub=None, numero=None, total=TOTAL):
    caixa(slide, MG, TOP_TIT, LARG - 2 * MG, 0.55, texto, tam=27, cor=AZUL,
          negrito=True)
    if sub:
        caixa(slide, MG, TOP_SUB, LARG - 2 * MG, 0.4, sub, tam=13, cor=TINTA2)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MG), Inches(Y_LINHA),
                                Inches(LARG - 2 * MG), Pt(1.1))
    ln.fill.solid()
    ln.fill.fore_color.rgb = RGBColor(0xBF, 0xD9, 0xEA)
    ln.line.fill.background()
    ln.shadow.inherit = False
    if numero is not None:
        rodape(slide, numero, total)


def rodape(slide, numero, total=TOTAL):
    caixa(slide, MG, Y_RODAPE, 7.0, 0.3,
          'Previsão de rebaixamento — Brasileirão Série A', tam=10, cor=TINTA2)
    caixa(slide, LARG - MG - 1.4, Y_RODAPE, 1.4, 0.3, f'{numero}/{total}',
          tam=10, cor=TINTA2, alinha=PP_ALIGN.RIGHT)


def tile(slide, x, y, w, h, numero, rotulo, cor=AZUL, tam_num=30, tam_rot=11):
    cx = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    cx.fill.solid()
    cx.fill.fore_color.rgb = FUNDOC
    cx.line.fill.background()
    cx.shadow.inherit = False
    cx.adjustments[0] = 0.10
    cx.text_frame.text = ''
    rico(slide, x, y + 0.13, w, h - 0.2,
         [[(numero, {'cor': cor, 'negrito': True, 'tam': tam_num})],
          [(rotulo, {'cor': TINTA2, 'tam': tam_rot})]],
         alinha=PP_ALIGN.CENTER, entrelinha=1.0, espaco_antes=3)
    return cx


def faixa(slide, x, y, w, h, texto, cor=AZUL, tam=15):
    mistura = RGBColor(*[int(c + (255 - c) * 0.88) for c in (cor[0], cor[1], cor[2])])
    cx = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    cx.fill.solid()
    cx.fill.fore_color.rgb = mistura
    cx.line.fill.background()
    cx.shadow.inherit = False
    cx.adjustments[0] = 0.14
    tf = cx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = texto
    r.font.name, r.font.size, r.font.color.rgb, r.font.bold = FONTE, Pt(tam), cor, True
    return cx


def figura(slide, nome, y, altura_max, centro_x=None, largura_max=None):
    """Insere figura escalada para caber, centralizada."""
    caminho = os.path.join(DIR_FIG, nome)
    with Image.open(caminho) as im:
        pw, ph = im.size
    prop = pw / ph
    h = altura_max
    w = h * prop
    lmax = largura_max if largura_max else (LARG - 2 * MG)
    if w > lmax:
        w = lmax
        h = w / prop
    cx = centro_x if centro_x else LARG / 2
    slide.shapes.add_picture(caminho, Inches(cx - w / 2), Inches(y),
                            Inches(w), Inches(h))
    return h


def legenda(slide, y, texto, tam=13.5):
    return rico(slide, MG, y, LARG - 2 * MG, 0.5, [[(texto, {})]], tam=tam,
                alinha=PP_ALIGN.CENTER)


_CONTA = [1]


def prox():
    """Numero do slide atual, incrementado a cada chamada."""
    _CONTA[0] += 1
    return _CONTA[0]


def novo():
    return prs.slides.add_slide(BRANCO_LAYOUT)


def secao(slide, rotulo):
    """Etiqueta discreta de seção no canto superior direito."""
    caixa(slide, LARG - MG - 3.6, TOP_TIT + 0.08, 3.6, 0.3, rotulo, tam=11,
          cor=VERM, negrito=True, alinha=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# 1. CAPA
# ════════════════════════════════════════════════════════════════════════════
s = novo()
brasao = os.path.join(RAIZ, 'tex', 'tcc_artigo', 'figuras', 'brasao_ufpb.png')
if os.path.exists(brasao):
    with Image.open(brasao) as im:
        pw, ph = im.size
    h = 1.05
    w = h * pw / ph
    s.shapes.add_picture(brasao, Inches(LARG / 2 - w / 2), Inches(0.55),
                         Inches(w), Inches(h))
rico(s, MG, 1.78, LARG - 2 * MG, 0.6,
     [[('UNIVERSIDADE FEDERAL DA PARAÍBA — CCSA', {'cor': TINTA2, 'tam': 12})],
      [('Bacharelado em Ciência de Dados para Negócios', {'cor': TINTA2, 'tam': 12})]],
     alinha=PP_ALIGN.CENTER, entrelinha=1.15, espaco_antes=2)
rico(s, MG, 2.72, LARG - 2 * MG, 1.4,
     [[('Previsão de Rebaixamento no', {'cor': AZUL, 'negrito': True, 'tam': 34})],
      [('Campeonato Brasileiro Série A', {'cor': AZUL, 'negrito': True, 'tam': 34})]],
     alinha=PP_ALIGN.CENTER, entrelinha=1.1, espaco_antes=0)
rico(s, MG, 4.18, LARG - 2 * MG, 0.8,
     [[('Uma comparação de modelos de ', {'cor': TINTA2, 'tam': 15}),
       ('Machine Learning', {'cor': TINTA2, 'tam': 15, 'italico': True}),
       (' com validação ', {'cor': TINTA2, 'tam': 15}),
       ('walk-forward', {'cor': TINTA2, 'tam': 15, 'italico': True})]],
     alinha=PP_ALIGN.CENTER)
rico(s, MG, 5.30, LARG - 2 * MG, 1.0,
     [[('Leonardo Feitosa Barroso', {'negrito': True, 'tam': 18})],
      [('Orientador: Prof. Dr. Hilton Martins de Brito Ramalho',
        {'cor': TINTA2, 'tam': 12})]],
     alinha=PP_ALIGN.CENTER, espaco_antes=6)
caixa(s, MG, 6.55, LARG - 2 * MG, 0.35, 'João Pessoa · 2026', tam=11, cor=TINTA2,
      alinha=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# 2. ROTEIRO
# ════════════════════════════════════════════════════════════════════════════
s = novo()
titulo(s, 'Roteiro', numero=prox())
itens_esq = [('1', 'O problema e a pergunta'), ('2', 'Dados e método'),
             ('3', 'Resultados')]
itens_dir = [('4', 'A prova real: 2025'), ('5', 'Considerações finais'),
             ('6', 'Demonstração do app')]
for col, itens in ((0, itens_esq), (1, itens_dir)):
    x = MG + 0.5 + col * 6.0
    pars = []
    for n, t in itens:
        cor = VERM if n == '6' else AZUL
        pars.append([(f'{n}   ', {'cor': cor, 'negrito': True, 'tam': 20}),
                     (t, {'negrito': n == '6', 'tam': 18})])
    rico(s, x, 2.35, 5.4, 2.4, pars, tam=18, entrelinha=1.3, espaco_antes=16)

# ════════════════════════════════════════════════════════════════════════════
# 3. O PROBLEMA
# ════════════════════════════════════════════════════════════════════════════
s = novo()
titulo(s, 'O problema',
       'Cair para a Série B é o pior resultado financeiro possível para um clube',
       numero=prox())
lt, gap = 3.35, 0.45
x0 = (LARG - (3 * lt + 2 * gap)) / 2
tile(s, x0, TOP_CONT, lt, 1.32, '4', 'clubes rebaixados por temporada', VERM)
tile(s, x0 + lt + gap, TOP_CONT, lt, 1.32, '20%',
     'das observações — classe minoritária', VERM)
tile(s, x0 + 2 * (lt + gap), TOP_CONT, lt, 1.32, '—',
     'receita de TV, patrocínio e elenco', AZUL)
marcadores(s, MG + 0.25, 3.42, LARG - 2 * MG - 0.5, 2.6, [
    [('A queda derruba receitas de transmissão, patrocínios e o valor do elenco.', {})],
    [('Os efeitos se arrastam por ', {}), ('várias temporadas', {'negrito': True}), ('.', {})],
    [('Antecipar o risco permite agir: reforços, renegociação, orçamento.', {})],
], tam=17, espaco=14)

# ════════════════════════════════════════════════════════════════════════════
# 4. A PERGUNTA
# ════════════════════════════════════════════════════════════════════════════
s = novo()
titulo(s, 'A pergunta', numero=prox())
rico(s, MG + 1.0, TOP_CONT + 0.20, LARG - 2 * MG - 2.0, 1.3,
     [[('Dá para estimar quem cai', {'cor': AZUL, 'negrito': True, 'tam': 30})],
      [('antes de a bola rolar', {'cor': AZUL, 'negrito': True, 'tam': 30,
                                  'italico': True}),
       ('?', {'cor': AZUL, 'negrito': True, 'tam': 30})]],
     alinha=PP_ALIGN.CENTER, entrelinha=1.15, espaco_antes=2)
marcadores(s, MG + 0.25, 3.62, LARG - 2 * MG - 0.5, 2.3, [
    [('A literatura foca em ', {}), ('resultado de partidas', {'negrito': True}),
     (' e em ', {}), ('ligas europeias', {'negrito': True}), ('.', {})],
    [('Não localizei trabalho sobre rebaixamento no Brasileirão usando ', {}),
     ('apenas informação pré-temporada', {'negrito': True}), ('.', {})],
    [('É justamente o cenário em que a previsão ', {}),
     ('serve para decidir algo', {'negrito': True}), ('.', {})],
], tam=17, espaco=13)
rico(s, MG + 0.25, 6.20, LARG - 2 * MG - 0.5, 0.5,
     [[('Busca no Google Scholar e Scopus: (relegation OR rebaixamento) AND '
        '(machine learning OR prediction).', {'cor': TINTA2, 'tam': 11})]])

# ════════════════════════════════════════════════════════════════════════════
# 5. OS DADOS
# ════════════════════════════════════════════════════════════════════════════
s = novo()
titulo(s, 'Os dados',
       'Transfermarkt, raspagem automatizada · dados públicos, uso acadêmico',
       numero=prox())
tile(s, x0, TOP_CONT, lt, 1.28, '220', 'observações rotuladas (2014–2024)', AZUL)
tile(s, x0 + lt + gap, TOP_CONT, lt, 1.28, '15', 'variáveis preditoras (features)', AZUL)
tile(s, x0 + 2 * (lt + gap), TOP_CONT, lt, 1.28, '2025',
     'temporada reservada para prever', VERM)
rico(s, MG + 0.3, 3.35, 5.6, 0.4,
     [[('3 de elenco', {'cor': AZUL, 'negrito': True, 'tam': 17}),
       ('  (o que o clube tem)', {'cor': TINTA2, 'tam': 14, 'italico': True})]])
marcadores(s, MG + 0.45, 3.88, 5.4, 1.8, [
    [('valor de mercado', {})], [('tamanho do plantel', {})],
    [('nº de estrangeiros', {})]], tam=15, espaco=7)
rico(s, MG + 6.6, 3.35, 6.0, 0.4,
     [[('12 de janela deslizante', {'cor': AZUL, 'negrito': True, 'tam': 17}),
       ('  (o que ele fez)', {'cor': TINTA2, 'tam': 14, 'italico': True})]])
marcadores(s, MG + 6.75, 3.88, 5.6, 1.8, [
    [('médias de 3 e 5 temporadas', {})],
    [('pontos, saldo, gols, vitórias, aproveitamento', {})]], tam=15, espaco=7)

# ════════════════════════════════════════════════════════════════════════════
# 6. TRATAMENTO DOS DADOS
# ════════════════════════════════════════════════════════════════════════════
s = novo()
titulo(s, 'Tratamento dos dados',
       'O dado bruto não vinha pronto — o trabalho de limpeza foi parte do método',
       numero=prox())
rico(s, MG + 0.3, TOP_CONT, 5.9, 0.4,
     [[('O problema dos nomes', {'cor': VERM, 'negrito': True, 'tam': 17})]])
rico(s, MG + 0.35, TOP_CONT + 0.5, 5.8, 1.5, [
    [('As duas páginas do Transfermarkt grafam o mesmo clube de formas '
      'diferentes — sem padronizar, o cruzamento falha:', {'tam': 13.5})],
], entrelinha=1.25)
# Exemplos de nomes divergentes -> padrão
exemplos = [('Clube Atletico Paranaense', 'Athletico Paranaense'),
            ('Athletico-PR', 'Athletico Paranaense'),
            ('CR Vasco da Gama', 'Vasco da Gama'),
            ('RB Bragantino', 'Bragantino'),
            ('Gremio FBPA', 'Gremio')]
yy = TOP_CONT + 1.32
for orig, pad in exemplos:
    rico(s, MG + 0.4, yy, 5.8, 0.3, [[
        (orig, {'cor': TINTA2, 'tam': 11.5, 'fonte': 'Consolas'}),
        ('  →  ', {'cor': VERM, 'tam': 11.5, 'negrito': True}),
        (pad, {'tam': 11.5, 'negrito': True, 'fonte': 'Consolas'}),
    ]], entrelinha=1.0)
    yy += 0.31
rico(s, MG + 0.35, yy + 0.12, 5.8, 0.6,
     [[('Um dicionário de ', {'cor': TINTA2, 'tam': 12}),
       ('normalização', {'cor': TINTA2, 'tam': 12, 'negrito': True}),
       (' com ~60 entradas + remoção de acentos resolveu o cruzamento por '
        'Clube + Temporada.', {'cor': TINTA2, 'tam': 12})]], entrelinha=1.25)

rico(s, MG + 6.9, TOP_CONT, 5.7, 0.4,
     [[('As outras etapas', {'cor': AZUL, 'negrito': True, 'tam': 17})]])
marcadores(s, MG + 6.95, TOP_CONT + 0.5, 5.6, 4.2, [
    [('Valores monetários', {'negrito': True}),
     (': remoção de “€”, “mi.”, “mil” e troca de vírgula por ponto.', {})],
    [('Alvo', {'negrito': True}),
     (': lista de rebaixados por temporada conferida ano a ano; '
      'o resto marcado como permanência.', {})],
    [('Auditoria do cruzamento', {'negrito': True}),
     (': checagem de registros com pontos ausentes para pegar merge que falhou.', {})],
    [('Valores faltantes', {'negrito': True}),
     (': recém-promovidos não têm histórico — imputação pela mediana '
      'do treino (15,5% das janelas).', {})],
    [('Escala', {'negrito': True}),
     (': padronização ajustada apenas no treino.', {})],
], tam=13.5, espaco=10)
faixa(s, MG + 0.3, 6.15, LARG - 2 * MG - 0.6, 0.6,
      '220 observações limpas e auditadas — 11 temporadas × 20 clubes, sem furo no cruzamento.',
      AZUL, tam=13.5)

# ════════════════════════════════════════════════════════════════════════════
# 7. O MÉTODO
# ════════════════════════════════════════════════════════════════════════════
s = novo()
titulo(s, 'O método',
       'Um fio condutor: nenhuma informação do futuro entra no modelo', numero=prox())
etapas = ['Coleta\nTransfermarkt', 'Janelas\nshift(1).rolling()',
          'Imputação e\npadronização\nsó no treino', 'Validação\nwalk-forward',
          'Ajuste de\nhiperparâm.', 'Teste\n2023–2024']
cw, cgap, ch = 1.83, 0.28, 1.30
cy = TOP_CONT + 0.55
cx0 = (LARG - (len(etapas) * cw + (len(etapas) - 1) * cgap)) / 2
for i, txt in enumerate(etapas):
    px = cx0 + i * (cw + cgap)
    cx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(px), Inches(cy),
                            Inches(cw), Inches(ch))
    cx.fill.solid()
    cx.fill.fore_color.rgb = RGBColor(0xE7, 0xF1, 0xF8)
    cx.line.color.rgb = RGBColor(0x9C, 0xC6, 0xE0)
    cx.line.width = Pt(1)
    cx.shadow.inherit = False
    cx.adjustments[0] = 0.12
    tf = cx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.05)
    for j, linha in enumerate(txt.split('\n')):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.05
        r = p.add_run()
        r.text = linha
        r.font.name, r.font.size, r.font.color.rgb = FONTE, Pt(11.5), TINTA
        r.font.bold = 'só no treino' in linha
    if i < len(etapas) - 1:
        seta = s.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, Inches(px + cw + 0.03),
            Inches(cy + ch / 2 - 0.055), Inches(cgap - 0.06), Inches(0.11))
        seta.fill.solid()
        seta.fill.fore_color.rgb = RGBColor(0x9C, 0xC6, 0xE0)
        seta.line.fill.background()
        seta.shadow.inherit = False
faixa(s, MG + 0.8, 4.72, LARG - 2 * MG - 1.6, 0.82,
      'As médias da temporada T usam só T−1 a T−w. Nunca a própria temporada prevista.',
      VERM, tam=15)

# ════════════════════════════════════════════════════════════════════════════
# 7. POR QUE WALK-FORWARD
# ════════════════════════════════════════════════════════════════════════════
s = novo()
titulo(s, 'Por que não a validação cruzada tradicional?', numero=prox())
lado = 0.30
esp = 0.075


def bloco(slide, x, y, cor):
    q = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                               Inches(lado), Inches(lado))
    q.fill.solid()
    q.fill.fore_color.rgb = cor
    q.line.fill.background()
    q.shadow.inherit = False
    q.adjustments[0] = 0.16
    return q


# esquerda: k-fold embaralhado
rico(s, MG + 0.3, TOP_CONT, 5.6, 0.4,
     [[('k-fold embaralhado', {'cor': ROXO, 'negrito': True, 'tam': 17})]])
ROXO_C = RGBColor(0xDF, 0xA9, 0xC6)
CINZA_C = RGBColor(0xC8, 0xCE, 0xD4)
for i in range(11):
    bloco(s, MG + 0.3 + i * (lado + esp), TOP_CONT + 0.55,
          ROXO_C if i % 2 == 0 else CINZA_C)
caixa(s, MG + 0.3, TOP_CONT + 1.02, 5.4, 0.3,
      '2014 → 2024, ordem embaralhada', tam=11, cor=TINTA2)
rico(s, MG + 0.3, TOP_CONT + 1.55, 5.4, 1.4,
     [[('Treina com temporadas ', {'tam': 15}),
       ('futuras', {'negrito': True, 'tam': 15}),
       (' para prever o passado.', {'tam': 15})],
      [('Otimista e irreal.', {'cor': ROXO, 'negrito': True, 'tam': 15})]],
     espaco_antes=12)

# direita: walk-forward
rico(s, MG + 6.75, TOP_CONT, 5.8, 0.4,
     [[('Walk-forward', {'cor': AZUL, 'negrito': True, 'tam': 17, 'italico': True}),
       (' (janela expansiva)', {'cor': AZUL, 'negrito': True, 'tam': 17})]])
AZUL_C = RGBColor(0x66, 0xAE, 0xD6)
VERM_C = RGBColor(0xE7, 0x93, 0x4D)
for r_ in range(5):
    n_tr = 6 + r_
    yy = TOP_CONT + 0.55 + r_ * (lado + esp)
    for i in range(n_tr):
        bloco(s, MG + 6.75 + i * (lado + esp), yy, AZUL_C)
    bloco(s, MG + 6.75 + n_tr * (lado + esp), yy, VERM_C)
caixa(s, MG + 6.75, TOP_CONT + 2.55, 5.6, 0.3, 'treino ■   validação ■',
      tam=11, cor=TINTA2)
rico(s, MG + 6.75, TOP_CONT + 2.95, 5.6, 0.8,
     [[('Sempre treina no passado e valida no ano seguinte.', {'tam': 15})]])

# ════════════════════════════════════════════════════════════════════════════
# 8. ESTABILIDADE (walk-forward)
# ════════════════════════════════════════════════════════════════════════════
s = novo()
titulo(s, 'Estabilidade ao longo das temporadas',
       'AUC-ROC em cada temporada de validação — 5 partições', numero=prox())
h = figura(s, 'slide_walkforward.png', TOP_CONT + 0.05, 4.45)
legenda(s, TOP_CONT + 0.05 + h + 0.14,
        'A Regressão Logística é a mais estável (±0,058); o LightGBM, o mais '
        'volátil (±0,151).')

# ════════════════════════════════════════════════════════════════════════════
# 9. TESTE INDEPENDENTE
# ════════════════════════════════════════════════════════════════════════════
s = novo()
titulo(s, 'Teste independente: 2023–2024',
       '40 observações nunca vistas em nenhuma etapa anterior', numero=prox())
figura(s, 'slide_roc.png', TOP_CONT + 0.05, 4.7, centro_x=4.35, largura_max=7.2)
rico(s, 8.55, TOP_CONT + 0.45, 4.2, 2.4, [
    [('■  ', {'cor': VERM, 'tam': 16}), ('LightGBM 0,877', {'negrito': True, 'tam': 16})],
    [('■  ', {'cor': VERDE, 'tam': 16}), ('Random Forest 0,844', {'tam': 16})],
    [('■  ', {'cor': AZUL, 'tam': 16}), ('Reg. Logística 0,828', {'tam': 16})],
    [('■  ', {'cor': ROXO, 'tam': 16}), ('XGBoost ', {'tam': 16}),
     ('0,652', {'negrito': True, 'tam': 16})],
], espaco_antes=11)
rico(s, 8.55, 4.62, 4.2, 1.5,
     [[('Mas o XGBoost tem a mesma acurácia que o LightGBM: 82,5%.',
        {'cor': VERM, 'negrito': True, 'tam': 17})]], entrelinha=1.25)

# ════════════════════════════════════════════════════════════════════════════
# 10. A ARMADILHA DA ACURÁCIA
# ════════════════════════════════════════════════════════════════════════════
s = novo()
titulo(s, 'A armadilha da acurácia',
       'Mesma acurácia, mesma matriz de confusão — AUC muito diferente', numero=prox())
rico(s, MG, TOP_CONT, LARG - 2 * MG, 0.4, [[
    ('■  ', {'cor': ROXO, 'tam': 15}),
    ('XGBoost 82,5% → AUC ', {'tam': 15}), ('0,652', {'negrito': True, 'tam': 15}),
    ('          ', {'tam': 15}),
    ('■  ', {'cor': VERM, 'tam': 15}),
    ('LightGBM 82,5% → AUC ', {'tam': 15}), ('0,877', {'negrito': True, 'tam': 15}),
]], alinha=PP_ALIGN.CENTER)
h = figura(s, 'slide_armadilha.png', TOP_CONT + 0.52, 3.55)
legenda(s, TOP_CONT + 0.52 + h + 0.16,
        'A acurácia só olha a decisão no limiar. O AUC olha a ordenação inteira '
        '— e é ela que serve à gestão.')

# ════════════════════════════════════════════════════════════════════════════
# 11. O QUE PESA
# ════════════════════════════════════════════════════════════════════════════
s = novo()
titulo(s, 'O que pesa no risco',
       'Coeficientes padronizados da Regressão Logística — 8 maiores em módulo',
       numero=prox())
h = figura(s, 'slide_importancia.png', TOP_CONT + 0.05, 4.25)
legenda(s, TOP_CONT + 0.05 + h + 0.14,
        'Valor de elenco é a variável dominante; vitórias recentes protegem.')
rico(s, MG, TOP_CONT + 0.05 + h + 0.58, LARG - 2 * MG, 0.4,
     [[('Sinais entre janelas da mesma métrica devem ser lidos em conjunto '
        '(multicolinearidade).', {'cor': TINTA2, 'tam': 11})]],
     alinha=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# 12. QUAL MODELO
# ════════════════════════════════════════════════════════════════════════════
s = novo()
titulo(s, 'Qual modelo levar para a prática?', numero=prox())
rico(s, MG + 0.3, TOP_CONT + 0.1, 5.6, 0.4,
     [[('LightGBM', {'cor': VERM, 'negrito': True, 'tam': 19})]])
marcadores(s, MG + 0.45, TOP_CONT + 0.68, 5.4, 2.2, [
    [('melhor AUC no teste: ', {}), ('0,877', {'negrito': True})],
    [('mas oscila entre temporadas: ±0,151', {})],
    [('caixa-preta', {})]], tam=16, cor_marca=VERM, espaco=10)
rico(s, MG + 6.75, TOP_CONT + 0.1, 5.8, 0.4,
     [[('Regressão Logística', {'cor': AZUL, 'negrito': True, 'tam': 19})]])
marcadores(s, MG + 6.9, TOP_CONT + 0.68, 5.5, 2.6, [
    [('AUC 0,828 — terceira', {})],
    [('mais estável', {'negrito': True}), (': ±0,058', {})],
    [('coeficientes interpretáveis', {})],
    [('probabilidades calibradas', {})]], tam=16, espaco=10)
faixa(s, MG + 0.6, 5.45, LARG - 2 * MG - 1.2, 0.95,
      'Escolhi a Regressão Logística: estabilidade e interpretabilidade acima do '
      'pico de desempenho.', AZUL, tam=16)

# ════════════════════════════════════════════════════════════════════════════
# 13. A PREVISÃO PARA 2025
# ════════════════════════════════════════════════════════════════════════════
s = novo()
titulo(s, 'A previsão para 2025',
       'Regressão Logística treinada em 2014–2024 · 10 maiores riscos', numero=prox())
h = figura(s, 'slide_previsao2025.png', TOP_CONT + 0.05, 4.5)
legenda(s, TOP_CONT + 0.05 + h + 0.16,
        'Juventude, Sport Recife, Vitória e Mirassol — os quatro maiores riscos.')

# ════════════════════════════════════════════════════════════════════════════
# 14. A PROVA REAL
# ════════════════════════════════════════════════════════════════════════════
s = novo()
titulo(s, 'A prova real: o campeonato aconteceu',
       'Validação retroativa contra o resultado oficial de 2025', numero=prox())
lt4, gap4 = 2.86, 0.30
x4 = (LARG - (4 * lt4 + 3 * gap4)) / 2
tile(s, x4, TOP_CONT, lt4, 1.25, '2 de 4', 'rebaixamentos previstos no corte',
     VERM, tam_num=26)
tile(s, x4 + lt4 + gap4, TOP_CONT, lt4, 1.25, '0,922',
     'AUC-ROC realizado na temporada', AZUL, tam_num=26)
tile(s, x4 + 2 * (lt4 + gap4), TOP_CONT, lt4, 1.25, '4 de 4',
     'rebaixados entre os 7 maiores riscos', AZUL, tam_num=26)
tile(s, x4 + 3 * (lt4 + gap4), TOP_CONT, lt4, 1.25, '59/64',
     'pares ordenados corretamente', AZUL, tam_num=26)
marcadores(s, MG + 0.25, 3.35, LARG - 2 * MG - 0.5, 3.0, [
    [('Melhor AUC de todos os períodos avaliados — acima do teste (0,828).', {})],
    [('O LightGBM, campeão no teste, caiu para ', {}), ('0,711', {'negrito': True}),
     (': a instabilidade vista no walk-forward se confirmou.', {})],
    [('Supera as heurísticas ingênuas — menor valor de elenco (0,906) e pior '
      'média de pontos (0,766).', {})],
], tam=16, espaco=13)

# ════════════════════════════════════════════════════════════════════════════
# 15. PREVISTO × ACONTECEU  (figura nova, ocupa o slide)
# ════════════════════════════════════════════════════════════════════════════
s = novo()
titulo(s, 'Previsto × aconteceu',
       'Ranking de risco previsto ligado à classificação final real de 2025',
       numero=prox())
figura(s, 'slide_previsto_real.png', TOP_CONT - 0.02, 5.20)

# ════════════════════════════════════════════════════════════════════════════
# 16. MIRASSOL
# ════════════════════════════════════════════════════════════════════════════
s = novo()
titulo(s, 'O erro que mais ensina: Mirassol',
       'Um caso atípico — e o tempo tem dado razão ao modelo', numero=prox())
tile(s, x0, TOP_CONT, lt, 1.25, '43,6%', 'risco previsto para 2025', VERM, tam_num=27)
tile(s, x0 + lt + gap, TOP_CONT, lt, 1.25, '4º', 'posição final em 2025', VERDE,
     tam_num=27)
tile(s, x0 + 2 * (lt + gap), TOP_CONT, lt, 1.25, '19º',
     'em 2026, campeonato em curso', VERM, tam_num=27)
marcadores(s, MG + 0.25, 3.28, LARG - 2 * MG - 0.5, 2.1, [
    [('Estreante sem histórico na Série A: as 12 features de janela receberam a '
      'mediana do treino.', {})],
    [('Em 2026 o clube aparece na ', {}),
     ('zona de rebaixamento', {'negrito': True}),
     (' — a fragilidade apontada não desapareceu, apenas demorou uma temporada.', {})],
], tam=16, espaco=12)
faixa(s, MG + 0.6, 5.42, LARG - 2 * MG - 1.2, 0.88,
      'Mirassol é outlier: sem a campanha atípica de 2025, teria caído.', AZUL,
      tam=16)
rico(s, MG, 6.42, LARG - 2 * MG, 0.4,
     [[('Classificação de 2026 em andamento (20 rodadas), consultada em jul./2026.',
        {'cor': TINTA2, 'tam': 10.5})]], alinha=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# 17. DEZ ANOS DE RISCO
# ════════════════════════════════════════════════════════════════════════════
s = novo()
titulo(s, 'Dez anos de risco em um quadro',
       'Modelo retreinado só com o passado de cada temporada · recorte de 16 clubes',
       numero=prox())
h = figura(s, 'slide_heatmap.png', TOP_CONT + 0.02, 4.55)
legenda(s, TOP_CONT + 0.02 + h + 0.14,
        'O “efeito elevador” no topo; os grandes em risco baixo e estável embaixo.')

# ════════════════════════════════════════════════════════════════════════════
# 18. CONSIDERAÇÕES FINAIS — retomada dos objetivos
# ════════════════════════════════════════════════════════════════════════════
s = novo()
titulo(s, 'Considerações finais: os objetivos',
       'Retomada dos cinco objetivos específicos', numero=prox())
secao(s, 'CONSIDERAÇÕES FINAIS')
objs = [
    ('i', 'Construir a base histórica 2014–2025 com elenco e desempenho',
     '220 observações rotuladas + 2025 para prever'),
    ('ii', 'Engenharia de features com janelas, sem vazamento',
     'shift(1).rolling(); imputação e escala só no treino'),
    ('iii', 'Treinar e otimizar quatro algoritmos com validação temporal',
     'walk-forward de 5 partições + RandomizedSearchCV'),
    ('iv', 'Comparar em teste independente com métrica adequada',
     'AUC-ROC em 2023–2024; LightGBM 0,877, Logística 0,828'),
    ('v', 'Aplicar o modelo final à previsão de 2025',
     'e validá-la contra o resultado oficial: AUC 0,922'),
]
yy = TOP_CONT + 0.05
for num, texto, detalhe in objs:
    rico(s, MG + 0.3, yy, LARG - 2 * MG - 0.6, 0.72, [
        [('✓  ', {'cor': VERDE, 'negrito': True, 'tam': 17}),
         (f'({num}) ', {'cor': TINTA2, 'negrito': True, 'tam': 15}),
         (texto, {'tam': 15.5, 'negrito': True})],
        [('        ', {'tam': 12}), (detalhe, {'cor': TINTA2, 'tam': 12.5})],
    ], entrelinha=1.1, espaco_antes=1)
    yy += 0.86
faixa(s, MG + 0.6, 6.10, LARG - 2 * MG - 1.2, 0.62,
      'O objetivo geral foi alcançado.', AZUL, tam=15)

# ════════════════════════════════════════════════════════════════════════════
# 19. CONSIDERAÇÕES FINAIS — contribuições
# ════════════════════════════════════════════════════════════════════════════
s = novo()
titulo(s, 'Considerações finais: contribuições', numero=prox())
secao(s, 'CONSIDERAÇÕES FINAIS')
contribs = [
    ('1', 'Pipeline reprodutível e sem vazamento temporal',
     'código, dados e experimentos em repositório público', AZUL),
    ('2', 'Demonstração empírica da armadilha da acurácia',
     'mesma acurácia e mesma matriz de confusão, AUC de 0,652 contra 0,877', AZUL),
    ('3', 'Histórico recente complementa o indicador financeiro',
     'elenco e desempenho carregam informação parcialmente distinta', AZUL),
    ('4', 'Uma previsão verificada contra a realidade',
     'não apenas prometida — 4 de 4 rebaixados entre os 7 maiores riscos', VERM),
]
yy = TOP_CONT + 0.22
for num, texto, detalhe, cor in contribs:
    rico(s, MG + 0.3, yy, LARG - 2 * MG - 0.6, 0.95, [
        [(f'{num}.  ', {'cor': cor, 'negrito': True, 'tam': 21}),
         (texto, {'tam': 17.5, 'negrito': True})],
        [('       ', {'tam': 13}), (detalhe, {'cor': TINTA2, 'tam': 13.5})],
    ], entrelinha=1.12, espaco_antes=2)
    yy += 1.17

# ════════════════════════════════════════════════════════════════════════════
# 20. CONSIDERAÇÕES FINAIS — limitações e trabalhos futuros
# ════════════════════════════════════════════════════════════════════════════
s = novo()
titulo(s, 'Considerações finais: limites e próximos passos', numero=prox())
secao(s, 'CONSIDERAÇÕES FINAIS')
rico(s, MG + 0.3, TOP_CONT + 0.05, 5.6, 0.4,
     [[('Limitações', {'cor': VERM, 'negrito': True, 'tam': 19})]])
marcadores(s, MG + 0.45, TOP_CONT + 0.62, 5.4, 3.0, [
    [('só dados pré-temporada: não capta lesão, troca de técnico, calendário', {})],
    [('~220 observações — amostra pequena aumenta a variância', {})],
    [('recém-promovidos dependem da imputação, como o Mirassol mostrou', {})],
], tam=15, cor_marca=VERM, espaco=11)
rico(s, MG + 6.75, TOP_CONT + 0.05, 5.8, 0.4,
     [[('Trabalhos futuros', {'cor': AZUL, 'negrito': True, 'tam': 19})]])
marcadores(s, MG + 6.9, TOP_CONT + 0.62, 5.5, 3.0, [
    [('retreino ', {}), ('durante', {'negrito': True}), (' a temporada', {})],
    [('features do mercado de transferências', {})],
    [('modelos sequenciais, com mais temporadas', {})],
    [('expansão para a Série B e outras ligas', {})],
], tam=15, espaco=11)
rico(s, MG + 0.3, 5.90, LARG - 2 * MG - 0.6, 0.8,
     [[('Análise de sensibilidade: imputação por média em vez de mediana → AUC '
        '0,828 (idêntico); janelas de 2 e 4 temporadas → 0,836. As decisões de '
        'projeto são robustas.', {'cor': TINTA2, 'tam': 12})]],
     alinha=PP_ALIGN.CENTER, entrelinha=1.2)

# ════════════════════════════════════════════════════════════════════════════
# 21. CONCLUSÃO → APP
# ════════════════════════════════════════════════════════════════════════════
s = novo()
titulo(s, 'Conclusão', numero=prox())
secao(s, 'CONSIDERAÇÕES FINAIS')
rico(s, MG + 0.9, TOP_CONT + 0.55, LARG - 2 * MG - 1.8, 1.6,
     [[('O modelo previu antes.', {'cor': AZUL, 'negrito': True, 'tam': 31})],
      [('O campeonato confirmou a ordenação.',
        {'cor': AZUL, 'negrito': True, 'tam': 31})]],
     alinha=PP_ALIGN.CENTER, entrelinha=1.18, espaco_antes=4)
faixa(s, MG + 1.4, 4.35, LARG - 2 * MG - 2.8, 0.95,
      'Agora, a ferramenta em funcionamento: demonstração do aplicativo', VERM,
      tam=17)
rico(s, MG, 5.85, LARG - 2 * MG, 0.5,
     [[('Código, dados e experimentos: github.com/leonardofeitos4/previsao-tcc',
        {'cor': TINTA2, 'tam': 12.5})]], alinha=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# 22. RESERVA
# ════════════════════════════════════════════════════════════════════════════
s = novo()
caixa(s, MG, TOP_TIT + 0.05, 8.0, 0.4, 'Reserva — números de apoio', tam=17,
      cor=TINTA2, negrito=True)
caixa(s, LARG - MG - 3.0, TOP_TIT + 0.1, 3.0, 0.3, '(não apresentar)', tam=11,
      cor=TINTA2, alinha=PP_ALIGN.RIGHT)
esq = [
    [('Walk-forward (médias)', {'negrito': True, 'tam': 13})],
    [('Reg. Logística 0,794 ± 0,058', {'tam': 12.5})],
    [('Random Forest 0,762 ± 0,071', {'tam': 12.5})],
    [('LightGBM 0,700 ± 0,151', {'tam': 12.5})],
    [('XGBoost 0,659 ± 0,085', {'tam': 12.5})],
    [('', {'tam': 6})],
    [('Fold 2020 (pandemia)', {'negrito': True, 'tam': 13})],
    [('LightGBM 0,438 · XGBoost 0,531 · RL 0,703', {'tam': 12.5})],
    [('', {'tam': 6})],
    [('Teste 2023–2024 (XGB e LGBM idênticos)', {'negrito': True, 'tam': 13})],
    [('3 acertos entre 8 rebaixados · 2 falsos positivos', {'tam': 12.5})],
    [('30 acertos de permanência', {'tam': 12.5})],
    [('', {'tam': 6})],
    [('Calibração (quartil de maior risco)', {'negrito': True, 'tam': 13})],
    [('previsto 63,6% · observado 50,0%', {'tam': 12.5})],
]
dir_ = [
    [('Base', {'negrito': True, 'tam': 13})],
    [('220 obs. rotuladas · 20% de prevalência', {'tam': 12.5})],
    [('15,5% de ausentes nas janelas', {'tam': 12.5})],
    [('Elenco: média €55,3M · DP €38,5M · máx. €214,2M', {'tam': 12.5})],
    [('', {'tam': 6})],
    [('Desbalanceamento', {'negrito': True, 'tam': 13})],
    [("class_weight='balanced' (RL, RF, LGBM)", {'tam': 12.5})],
    [('scale_pos_weight=4 (XGBoost)', {'tam': 12.5})],
    [('', {'tam': 6})],
    [('Otimização', {'negrito': True, 'tam': 13})],
    [('RandomizedSearchCV, 30 candidatos,', {'tam': 12.5})],
    [('TimeSeriesSplit de 5 folds, critério AUC', {'tam': 12.5})],
    [('', {'tam': 6})],
    [('Rebaixados de 2025', {'negrito': True, 'tam': 13})],
    [('Sport (20º) · Juventude (19º)', {'tam': 12.5})],
    [('Fortaleza (18º) · Ceará (17º)', {'tam': 12.5})],
]
rico(s, MG + 0.2, 1.35, 6.0, 5.4, esq, entrelinha=1.18, espaco_antes=2)
rico(s, MG + 6.7, 1.35, 6.0, 5.4, dir_, entrelinha=1.18, espaco_antes=2)

# ════════════════════════════════════════════════════════════════════════════
prs.save(SAIDA)
print(f'[ok] {os.path.relpath(SAIDA, RAIZ)}  —  {len(prs.slides)} slides')
